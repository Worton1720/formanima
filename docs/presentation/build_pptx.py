#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Генератор презентации защиты ВКР FORMANIMA (академический шаблон ИС).
12 слайдов: титул · цель/задачи · актуальность+НПБ · предметная область ·
AS-IS+НПБ · проблема · требования · Use Case+ER · TO-BE+архитектура ·
движок геймификации · аналоги · благодарность.

Запуск:  uv run --with python-pptx python build_pptx.py
Результат: defense.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "..", "screenshots")
TOTAL = 12

# --- "Кузнечная" тёмная палитра (в тон UI приложения) ---
OBSIDIAN = RGBColor(0x14, 0x12, 0x10)   # фон
PANEL    = RGBColor(0x1F, 0x1B, 0x17)   # подложки
GOLD     = RGBColor(0xE8, 0xB4, 0x4A)   # заголовки/акцент
EMBER    = RGBColor(0xD9, 0x6E, 0x2E)   # вторичный акцент
TEXT     = RGBColor(0xEC, 0xE6, 0xDC)   # основной текст
MUTED    = RGBColor(0xA9, 0x9E, 0x8E)   # приглушённый
DANGER   = RGBColor(0xCF, 0x5B, 0x4E)   # проблема/«минус»
OK       = RGBColor(0x7F, 0xB8, 0x6B)   # «плюс»

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- хелперы
def bg(slide, color=OBSIDIAN):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(r, text, size, color, bold=False, italic=False, font="Calibri"):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def accent_bar(slide, t=Inches(1.32), w=Inches(2.0)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), t, w, Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background(); bar.shadow.inherit = False


def title(slide, text, kicker=None):
    tb, tf = box(slide, Inches(0.9), Inches(0.45), Inches(11.6), Inches(0.95))
    set_run(tf.paragraphs[0].add_run(), text, 30, GOLD, bold=True)
    accent_bar(slide)
    if kicker:
        kb, ktf = box(slide, Inches(0.92), Inches(0.05), Inches(11.5), Inches(0.4))
        set_run(ktf.paragraphs[0].add_run(), kicker, 13, EMBER, bold=True)


def bullets(slide, items, left=Inches(0.95), top=Inches(1.7),
            width=Inches(11.4), height=Inches(5.3), size=19, gap=10):
    tb, tf = box(slide, left, top, width, height)
    first = True
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        marker = "▸ " if lvl == 0 else "•  "
        set_run(p.add_run(), marker, size, GOLD if lvl == 0 else EMBER, bold=True)
        set_run(p.add_run(), txt, size if lvl == 0 else size - 2,
                TEXT if lvl == 0 else MUTED, bold=False)
    return tb


def footer(slide, n):
    fb, ftf = box(slide, Inches(11.4), Inches(7.0), Inches(1.8), Inches(0.4))
    p = ftf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), f"FORMANIMA · {n}/{TOTAL}", 10, MUTED)


def rrect(slide, x, y, w, h, fill=PANEL, line=GOLD, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def shape_text(sh, text, size, color, bold=True):
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, size, color, bold=bold)


def line(slide, x1, y1, x2, y2, color=MUTED, width=1.25):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    cn.shadow.inherit = False
    return cn


def label(slide, x, y, w, h, text, size, color, align=PP_ALIGN.CENTER, bold=False, italic=False):
    _, tf = box(slide, x, y, w, h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    set_run(p.add_run(), text, size, color, bold=bold, italic=italic)


def actor(slide, cx, top, name):
    """Простой человечек: голова (овал) + тело (треугольник) + подпись."""
    head = rrect(slide, cx - Inches(0.16), top, Inches(0.32), Inches(0.32),
                 fill=GOLD, line=None, shape=MSO_SHAPE.OVAL)
    body = rrect(slide, cx - Inches(0.26), top + Inches(0.34), Inches(0.52), Inches(0.5),
                 fill=GOLD, line=None, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
    label(slide, cx - Inches(0.95), top + Inches(0.86), Inches(1.9), Inches(0.4),
          name, 13, TEXT, bold=True)


def idef0(slide, x, y, w, h, proc, control, inp, out, mech,
          note=None, note_color=MUTED, proc_color=GOLD):
    """Упрощённая нотация IDEF0: вход(слева)→процесс→выход(справа),
    управление(сверху), механизм(снизу)."""
    cx = Emu(int(x + w / 2)); cy = Emu(int(y + h / 2))
    box_proc = rrect(slide, x, y, w, h, fill=PANEL, line=proc_color, lw=1.75)
    shape_text(box_proc, proc, 15, proc_color, bold=True)
    arm = Inches(0.95)
    # вход (слева → внутрь)
    line(slide, Emu(int(x - arm)), cy, x, cy, color=EMBER, width=2)
    label(slide, Emu(int(x - arm - Inches(1.7))), Emu(int(cy - Inches(0.55))),
          Inches(1.75), Inches(0.9), inp, 11, MUTED, align=PP_ALIGN.RIGHT)
    # выход (изнутри → вправо)
    line(slide, Emu(int(x + w)), cy, Emu(int(x + w + arm)), cy, color=OK, width=2)
    label(slide, Emu(int(x + w + Inches(0.05))), Emu(int(cy - Inches(0.55))),
          Inches(1.85), Inches(0.9), out, 11, MUTED, align=PP_ALIGN.LEFT)
    # управление (сверху → вниз)
    line(slide, cx, Emu(int(y - arm)), cx, y, color=GOLD, width=2)
    label(slide, Emu(int(cx - Inches(1.6))), Emu(int(y - arm - Inches(0.45))),
          Inches(3.2), Inches(0.55), control, 11, MUTED)
    # механизм (снизу → вверх)
    line(slide, cx, Emu(int(y + h + arm)), cx, Emu(int(y + h)), color=GOLD, width=2)
    label(slide, Emu(int(cx - Inches(1.7))), Emu(int(y + h + Inches(0.05))),
          Inches(3.4), Inches(0.7), mech, 11, MUTED)
    if note:
        label(slide, x, Emu(int(y + h + arm + Inches(0.55))), w, Inches(0.8),
              note, 12, note_color, bold=True)


# ============================================================ 1 — Титул
s = prs.slides.add_slide(BLANK); bg(s)
strip = rrect(s, 0, Inches(2.9), SW, Pt(3), fill=GOLD, line=None, shape=MSO_SHAPE.RECTANGLE)
label(s, Inches(1.0), Inches(0.7), Inches(11.3), Inches(0.6),
      "Выпускная квалификационная работа", 16, EMBER, align=PP_ALIGN.LEFT, bold=True)
_, tf = box(s, Inches(1.0), Inches(1.35), Inches(11.3), Inches(1.5))
set_run(tf.paragraphs[0].add_run(),
        "Разработка веб-приложения для трекинга привычек и целей "
        "с элементами геймификации", 30, TEXT, bold=True)
label(s, Inches(1.0), Inches(3.15), Inches(11.3), Inches(0.7),
      "Проект FORMANIMA  ·  Vue 3 + NestJS + Prisma + PostgreSQL",
      20, GOLD, align=PP_ALIGN.LEFT, bold=True)
_, tf = box(s, Inches(1.0), Inches(5.35), Inches(11.3), Inches(1.7))
for i, ln in enumerate([
        "Выполнил: ____________________,  группа ________",
        "Научный руководитель: ____________________",
        "____________________________________  (вуз)",
        "2026"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    set_run(p.add_run(), ln, 15, MUTED)

# ============================================================ 2 — Цель и задачи
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Цель и задачи работы")
bullets(s, [
    "Для кого: пользователи, формирующие полезные привычки и достигающие личных целей — студенты, специалисты, все, кто следит за продуктивностью и здоровьем",
    "Цель: разработать веб-приложение для трекинга привычек и целей с элементами геймификации, мотивирующее пользователя регулярно выполнять свои цели",
    "Задачи:",
    ("проанализировать предметную область и существующие аналоги", 1),
    ("спроектировать архитектуру системы и базу данных", 1),
    ("реализовать модули: клиент (SPA), серверный REST API, движок геймификации", 1),
    ("протестировать систему и оценить результат", 1),
], top=Inches(1.6), gap=9)
footer(s, 2)

# ============================================================ 3 — Актуальность + НПБ-хук
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Актуальность темы", "Зачем это нужно сейчас")
bullets(s, [
    "Цифровое благополучие и продуктивность — растущий тренд: люди массово ставят цели и заводят полезные привычки",
    "Главная проблема — удержание: по исследованию Ф. Лалли (2010) привычка закрепляется в среднем за 66 дней, а большинство бросает раньше",
    "Приложение хранит персональные данные (e-mail, пароль) → защита обязательна по ФЗ-152 «О персональных данных»",
    "Решение — геймификация на открытом стеке (курс на импортозамещение)",
], top=Inches(1.55), width=Inches(7.3), height=Inches(5.2), size=18, gap=14)
# Столбчатая диаграмма: 21 / 66 / 254 дня
bars = [("21", 21, "миф", MUTED), ("66", 66, "реальность", GOLD), ("254", 254, "максимум", EMBER)]
bx = Inches(8.6); base_y = Inches(6.1); max_h = Inches(3.2); bw = Inches(0.95); step = Inches(1.25)
label(s, Inches(8.4), Inches(1.6), Inches(4.4), Inches(0.7),
      "Дней до автоматизма привычки\n(Ф. Лалли, 2010)", 13, MUTED, bold=True)
for i, (val, days, cap, col) in enumerate(bars):
    h = Emu(int(max_h * days / 254))
    x = Emu(int(bx + i * step))
    y = Emu(int(base_y - h))
    rrect(s, x, y, bw, h, fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
    label(s, Emu(int(x - Inches(0.25))), Emu(int(y - Inches(0.45))), Emu(int(bw + Inches(0.5))),
          Inches(0.4), val, 15, col, bold=True)
    label(s, Emu(int(x - Inches(0.25))), Emu(int(base_y + Inches(0.05))),
          Emu(int(bw + Inches(0.5))), Inches(0.4), cap, 11, MUTED)
footer(s, 3)

# ============================================================ 4 — Предметная область ч.1
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Предметная область: сущности и процессы")
bullets(s, [
    "Область — формирование и удержание полезных привычек и достижение личных целей пользователя",
    "Участники: Пользователь (ведёт привычки и цели) и Администратор (управляет пользователями)",
    "Базовые процессы: регистрация → создание привычки или цели → ежедневная отметка выполнения → автоматический пересчёт прогресса → просмотр статистики",
], top=Inches(1.55), height=Inches(2.6), size=18, gap=14)
# Чипы сущностей
chips = ["Пользователь", "Привычка", "Действие", "Цель", "Отметка", "Достижение", "Игровой профиль"]
label(s, Inches(0.95), Inches(4.35), Inches(11.0), Inches(0.4),
      "Ключевые сущности системы:", 14, GOLD, align=PP_ALIGN.LEFT, bold=True)
cx = Inches(0.95); cy = Inches(4.95); chh = Inches(0.7); pad = Inches(0.28); gap = Inches(0.25)
for ch in chips:
    cw = Emu(int(Inches(0.12) * len(ch) + pad * 2))
    if cx + cw > SW - Inches(0.9):
        cx = Inches(0.95); cy = Emu(int(cy + chh + Inches(0.2)))
    chip = rrect(s, cx, cy, cw, chh, fill=PANEL, line=EMBER, lw=1)
    shape_text(chip, ch, 13, TEXT, bold=True)
    cx = Emu(int(cx + cw + gap))
footer(s, 4)

# ============================================================ 5 — AS-IS + НПБ
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Текущий процесс (AS-IS) и нормативная база")
idef0(s, Inches(2.7), Inches(2.9), Inches(3.6), Inches(1.4),
      "Ведение привычек\nвручную",
      control="дисциплина, память",
      inp="желание,\nцель",
      out="несистемный\nрезультат,\nсрывы",
      mech="блокнот / заметки в телефоне",
      note="Минусы: нет автоподсчёта · легко забыть · нет мотивации · нет статистики",
      note_color=DANGER, proc_color=DANGER)
# НПБ справа
label(s, Inches(8.4), Inches(1.75), Inches(4.4), Inches(0.4),
      "Нормативно-правовая база", 16, GOLD, align=PP_ALIGN.LEFT, bold=True)
bullets(s, [
    "ФЗ-152 «О персональных данных»",
    "ФЗ-149 «Об информации, ИТ и о защите информации»",
    "ГОСТ 34.601-90 — стадии создания АС",
    "ГОСТ 34.602-89 — техническое задание на АС",
    "Отраслевые: REST, JWT (RFC 7519), OWASP",
], left=Inches(8.4), top=Inches(2.3), width=Inches(4.5), height=Inches(4.4),
   size=14, gap=14)
footer(s, 5)

# ============================================================ 6 — Проблема ПО
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Проблема предметной области")
bullets(s, [
    "Низкое удержание: большинство бросает новые привычки за 2–3 недели — нет внешней мотивации",
    "Ручной учёт ненадёжен: легко забыть отметить, теряется история, нет напоминаний",
    "Нет наглядного прогресса: пользователь не видит динамику и теряет интерес",
    "Разрозненность: привычки, цели, финансы, питание ведутся в разных приложениях",
], top=Inches(1.6), height=Inches(3.6), size=19, gap=16)
# Вывод
concl = rrect(s, Inches(0.95), Inches(5.5), Inches(11.4), Inches(1.2), fill=PANEL, line=GOLD, lw=1.5)
tf = concl.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Вывод: ", 18, GOLD, bold=True)
set_run(p.add_run(),
        "ручное ведение и разрозненные трекеры не обеспечивают долгую мотивацию → "
        "нужна единая система с геймификацией", 18, TEXT)
footer(s, 6)

# ============================================================ 7 — Требования
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Требования к системе", "Постановка задачи")
# два столбца
colA = rrect(s, Inches(0.9), Inches(1.65), Inches(5.85), Inches(5.1), fill=PANEL, line=GOLD, lw=1)
colB = rrect(s, Inches(7.0), Inches(1.65), Inches(5.45), Inches(5.1), fill=PANEL, line=EMBER, lw=1)
label(s, Inches(1.1), Inches(1.8), Inches(5.4), Inches(0.5),
      "Функциональные", 18, GOLD, align=PP_ALIGN.LEFT, bold=True)
bullets(s, [
    "Регистрация и авторизация по ролям (JWT)",
    "Управление привычками, действиями и целями (CRUD)",
    "Ежедневная отметка выполнения",
    "Авто-пересчёт опыта, уровней, рангов, достижений",
    "Статистика: графики, тепловая карта, стрики",
    "Доп. модули: финансы и калории (КБЖУ)",
    "Админ-панель: управление пользователями",
], left=Inches(1.15), top=Inches(2.4), width=Inches(5.35), height=Inches(4.2),
   size=14, gap=9)
label(s, Inches(7.2), Inches(1.8), Inches(5.0), Inches(0.5),
      "Нефункциональные", 18, EMBER, align=PP_ALIGN.LEFT, bold=True)
bullets(s, [
    "Безопасность: хеширование bcrypt, JWT, защита ПДн (ФЗ-152)",
    "Веб-интерфейс (SPA), работа в браузере, адаптивность",
    "Документированный REST API (Swagger, /api/docs)",
    "Модульность и расширяемость (8 модулей NestJS)",
    "Открытый стек, запуск в Docker (импортонезависимость)",
], left=Inches(7.25), top=Inches(2.4), width=Inches(4.95), height=Inches(4.2),
   size=14, gap=11)
footer(s, 7)

# ============================================================ 8 — Концепция: Use Case + ER
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Концепция ИС: Use Case и модель данных")
# --- слева: диаграмма вариантов использования ---
label(s, Inches(0.9), Inches(1.55), Inches(5.8), Inches(0.4),
      "Диаграмма вариантов использования", 14, GOLD, align=PP_ALIGN.LEFT, bold=True)
usecases = [
    ("Вход и регистрация", Inches(2.35)),
    ("Вести привычки и цели", Inches(3.15)),
    ("Отмечать выполнение", Inches(3.95)),
    ("Смотреть статистику\nи достижения", Inches(4.85)),
]
uc_x = Inches(2.55); uc_w = Inches(2.7); uc_h = Inches(0.62)
actor(s, Inches(1.25), Inches(2.9), "Пользователь")
for txt, uy in usecases:
    oval = rrect(s, uc_x, uy, uc_w, uc_h, fill=PANEL, line=GOLD, lw=1, shape=MSO_SHAPE.OVAL)
    shape_text(oval, txt, 11, TEXT, bold=False)
    line(s, Inches(1.55), Inches(3.25), uc_x, Emu(int(uy + uc_h / 2)), color=MUTED, width=1)
# админ
admin_uc = rrect(s, uc_x, Inches(5.75), uc_w, uc_h, fill=PANEL, line=EMBER, lw=1, shape=MSO_SHAPE.OVAL)
shape_text(admin_uc, "Управление\nпользователями", 11, TEXT, bold=False)
actor(s, Inches(6.05), Inches(5.55), "Админ")
line(s, Inches(5.75), Inches(6.05), Emu(int(uc_x + uc_w)), Emu(int(Inches(5.75) + uc_h / 2)),
     color=MUTED, width=1)
# --- справа: ER-диаграмма ---
label(s, Inches(7.0), Inches(1.55), Inches(5.5), Inches(0.4),
      "Модель данных (ER), 13 сущностей", 14, GOLD, align=PP_ALIGN.LEFT, bold=True)
# центральная сущность User
user_box = rrect(s, Inches(9.0), Inches(3.55), Inches(1.7), Inches(0.7), fill=EMBER, line=GOLD, lw=1.5)
shape_text(user_box, "User", 14, OBSIDIAN, bold=True)
ucx, ucy = Inches(9.85), Inches(3.9)
satellites = [
    ("UserRank\n(опыт/уровень)", Inches(7.05), Inches(2.1), "1—1"),
    ("Habit → Action", Inches(7.05), Inches(3.7), "1—∞"),
    ("Goal → Progress", Inches(7.05), Inches(5.2), "1—∞"),
    ("Achievement", Inches(11.0), Inches(2.1), "1—∞"),
    ("Transaction\n/ Budget", Inches(11.0), Inches(3.7), "1—∞"),
    ("FoodEntry", Inches(11.0), Inches(5.2), "1—∞"),
]
ebw, ebh = Inches(1.95), Inches(0.72)
for txt, ex, ey, card in satellites:
    eb = rrect(s, ex, ey, ebw, ebh, fill=PANEL, line=MUTED, lw=1)
    shape_text(eb, txt, 11, TEXT, bold=False)
    line(s, ucx, ucy, Emu(int(ex + ebw / 2)), Emu(int(ey + ebh / 2)), color=MUTED, width=1)
    label(s, Emu(int((ucx + ex + ebw / 2) / 2) - int(Inches(0.4))),
          Emu(int((ucy + ey + ebh / 2) / 2) - int(Inches(0.2))),
          Inches(0.8), Inches(0.3), card, 10, GOLD, bold=True)
footer(s, 8)

# ============================================================ 9 — Описание ИС: TO-BE + архитектура
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Описание ИС: автоматизированный процесс (TO-BE)")
idef0(s, Inches(4.85), Inches(2.9), Inches(3.6), Inches(1.15),
      "Трекинг привычек\nс геймификацией",
      control="правила геймификации, ФЗ-152, JWT",
      inp="отметка\nвыполнения",
      out="прогресс,\nуровень,\nстатистика",
      mech="браузер + Vue + NestJS + PostgreSQL",
      proc_color=OK)
# поток данных снизу
flow = ["Vue 3 + Pinia", "axios", "NestJS\n+ Guards", "Prisma ORM", "PostgreSQL"]
n = len(flow); gp = Inches(0.3)
total_w = SW - Inches(1.8) - gp * (n - 1)
bw = Emu(int(total_w / n)); x = Inches(0.9); y = Inches(6.0); bh = Inches(0.9)
for i, lab in enumerate(flow):
    sh = rrect(s, x, y, bw, bh, fill=PANEL, line=GOLD, lw=1.25)
    shape_text(sh, lab, 13, TEXT, bold=True)
    if i < n - 1:
        label(s, Emu(int(x + bw)), y, gp, bh, "→", 20, EMBER, bold=True)
    x = Emu(int(x + bw + gp))
label(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.4),
      "Архитектура (монорепозиторий, Docker):", 14, GOLD, align=PP_ALIGN.LEFT, bold=True)
footer(s, 9)

# ============================================================ 10 — Ядро: движок геймификации
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Ядро системы: движок геймификации", "Ключевой модуль")
bullets(s, [
    "При каждой отметке выполнения система сама пересчитывает игровые метрики (recalculate)",
    "Опыт: страйки × 10 + идеальные дни × 50 + награды достижений",
    "Уровень = ⌊√(XP / 100)⌋ + 1 — каждый уровень даётся труднее",
    "Ранги: подмастерье → ремесленник → мастер → гроссмейстер",
    "12 достижений: стрики, идеальная неделя, питание, финансы",
], top=Inches(1.55), width=Inches(6.2), height=Inches(3.6), size=15, gap=10)
# таблица уровней
rows = [("Уровень", "Нужно XP"), ("2", "100"), ("3", "400"), ("4", "900"), ("5", "1600")]
gt = s.shapes.add_table(len(rows), 2, Inches(0.95), Inches(5.1), Inches(3.0), Inches(1.9)).table
gt.columns[0].width = Inches(1.5); gt.columns[1].width = Inches(1.5)
for ri, (a, b) in enumerate(rows):
    for ci, val in enumerate((a, b)):
        cell = gt.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = EMBER if ri == 0 else PANEL
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        set_run(para.add_run(), val, 13, OBSIDIAN if ri == 0 else TEXT, bold=(ri == 0))
# скриншот дашборда
shot = os.path.join(SHOTS, "04-dashboard.png")
if os.path.exists(shot):
    pic = s.shapes.add_picture(shot, Inches(7.4), Inches(1.7), width=Inches(5.4))
    pic.line.color.rgb = GOLD; pic.line.width = Pt(1)
    label(s, Inches(7.4), Inches(5.7), Inches(5.4), Inches(0.4),
          "Дашборд: ранг, уровень, опыт, стрики", 12, MUTED, italic=True)
footer(s, 10)

# ============================================================ 11 — Аналоги
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "Сравнение с аналогами", "Почему наше решение")
crit = [
    ("Возможность", "FORMANIMA", "Habitica", "Forest", "Loop"),
    ("Привычки и действия", "✓", "✓", "—", "✓"),
    ("Цели с этапами", "✓", "~", "—", "—"),
    ("Геймификация (уровни/ранги)", "✓", "✓", "~", "—"),
    ("Учёт финансов", "✓", "—", "—", "—"),
    ("Учёт калорий (КБЖУ)", "✓", "—", "—", "—"),
    ("Всё в одном приложении", "✓", "—", "—", "—"),
    ("Русский язык / свой хостинг", "✓", "~", "—", "~"),
]
nrows = len(crit); ncols = 5
t = s.shapes.add_table(nrows, ncols, Inches(0.9), Inches(1.6), Inches(11.5), Inches(4.4)).table
t.columns[0].width = Inches(4.3)
for c in range(1, ncols):
    t.columns[c].width = Inches(1.8)
for ri, row in enumerate(crit):
    for ci, val in enumerate(row):
        cell = t.cell(ri, ci)
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = EMBER if ci == 1 else PANEL
            col = OBSIDIAN if ci == 1 else GOLD
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ci != 1 else RGBColor(0x2A, 0x23, 0x18)
            col = OK if (ci == 1 and val == "✓") else (TEXT if ci == 0 else MUTED)
            if ci == 1 and val == "✓":
                col = OK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        set_run(para.add_run(), val, 13, col, bold=(ri == 0 or ci == 1))
label(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.7),
      "Ниша FORMANIMA — единственное решение, объединяющее привычки, цели, "
      "финансы и питание под одной системой мотивации.", 14, GOLD, bold=True)
footer(s, 11)

# ============================================================ 12 — Благодарность
s = prs.slides.add_slide(BLANK); bg(s)
rrect(s, 0, Inches(3.5), SW, Pt(3), fill=GOLD, line=None, shape=MSO_SHAPE.RECTANGLE)
label(s, Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.2),
      "Доклад окончен, благодарю за внимание!", 30, GOLD, bold=True)
label(s, Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.6),
      "FORMANIMA · 2026", 16, MUTED)

out = os.path.join(HERE, "defense.pptx")
prs.save(out)
print("OK saved:", out, "| slides:", len(prs.slides._sldIdLst))
